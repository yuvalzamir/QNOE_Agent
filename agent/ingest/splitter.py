"""Chunking strategies for different content types.

Supported:
  - Python files (.py)        — function/class-level via AST
  - Jupyter notebooks (.ipynb)— cell-level (code + markdown)
  - Markdown / text (.md, .txt, .rst) — paragraph-level
  - PDFs                      — section-level via Docling (if installed)

Each chunk is a dict:
  {
      "text":       str,
      "source":     str,   # relative file path
      "repo":       str,
      "chunk_type": "prose" | "code",
      "start_line": int | None,
  }
"""
import ast
import json
import logging
import os
from pathlib import Path

logger = logging.getLogger(__name__)

MAX_CHUNK_CHARS = 2000
OVERLAP_CHARS = 200


# ── Python ─────────────────────────────────────────────────────────────────────

def _chunk_python(path: Path, repo: str) -> list[dict]:
    source = path.read_text(encoding="utf-8", errors="replace")
    try:
        tree = ast.parse(source)
    except SyntaxError:
        # Fall back to plain text split
        return _chunk_text(path, repo, chunk_type="code")

    chunks: list[dict] = []
    lines = source.splitlines()

    def extract(node) -> None:
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
            start = node.lineno - 1
            end = node.end_lineno
            text = "\n".join(lines[start:end])
            if len(text) > MAX_CHUNK_CHARS:
                # Split oversized functions into sliding windows
                for chunk_text in _sliding_window(text):
                    chunks.append({
                        "text": chunk_text,
                        "source": str(path),
                        "repo": repo,
                        "chunk_type": "code",
                        "start_line": start + 1,
                    })
            else:
                chunks.append({
                    "text": text,
                    "source": str(path),
                    "repo": repo,
                    "chunk_type": "code",
                    "start_line": start + 1,
                })
        for child in ast.iter_child_nodes(node):
            extract(child)

    extract(tree)

    # If no functions/classes found, chunk as plain text
    if not chunks:
        return _chunk_text(path, repo, chunk_type="code")

    return chunks


# ── Jupyter notebooks ──────────────────────────────────────────────────────────

def _chunk_notebook(path: Path, repo: str) -> list[dict]:
    try:
        nb = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    except json.JSONDecodeError:
        logger.warning("Could not parse notebook: %s", path)
        return []

    chunks: list[dict] = []
    for i, cell in enumerate(nb.get("cells", [])):
        cell_type = cell.get("cell_type", "")
        source_lines = cell.get("source", [])
        text = "".join(source_lines).strip()
        if not text:
            continue

        # Check for error outputs
        has_error = any(
            out.get("output_type") == "error"
            for out in cell.get("outputs", [])
        )
        if has_error:
            text = "[ERROR OUTPUT]\n" + text

        chunk_type = "code" if cell_type == "code" else "prose"
        for part in _sliding_window(text):
            chunks.append({
                "text": part,
                "source": str(path),
                "repo": repo,
                "chunk_type": chunk_type,
                "start_line": i,
            })

    return chunks


# ── Markdown / plain text ──────────────────────────────────────────────────────

def _chunk_text(path: Path, repo: str, chunk_type: str = "prose") -> list[dict]:
    text = path.read_text(encoding="utf-8", errors="replace")
    chunks: list[dict] = []
    # Split on blank lines (paragraph boundaries)
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    buffer = ""
    for para in paragraphs:
        if len(buffer) + len(para) > MAX_CHUNK_CHARS and buffer:
            chunks.append({
                "text": buffer.strip(),
                "source": str(path),
                "repo": repo,
                "chunk_type": chunk_type,
                "start_line": None,
            })
            buffer = buffer[-OVERLAP_CHARS:] + "\n\n" + para
        else:
            buffer = (buffer + "\n\n" + para) if buffer else para
    if buffer.strip():
        chunks.append({
            "text": buffer.strip(),
            "source": str(path),
            "repo": repo,
            "chunk_type": chunk_type,
            "start_line": None,
        })
    return chunks


def _chunk_pptx(path: Path, repo: str) -> list[dict]:
    """Extract text slide-by-slide from a PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed — skipping PPTX: %s", path)
        return []
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        logger.warning("Could not open PPTX %s: %s", path, exc)
        return []
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        text = "\n".join(parts)
        if not text.strip():
            continue
        for part in _sliding_window(f"[Slide {i}] " + text):
            chunks.append({"text": part, "source": str(path), "repo": repo, "chunk_type": "prose", "start_line": i})
    return chunks


def _chunk_docx(path: Path, repo: str) -> list[dict]:
    """Extract text from a Word document — paragraphs and tables in document order."""
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        logger.warning("python-docx not installed — skipping DOCX: %s", path)
        return []
    try:
        doc = Document(str(path))
    except Exception as exc:
        logger.warning("Could not open DOCX %s: %s", path, exc)
        return []

    def _iter_blocks(doc):
        """Yield paragraphs and tables in document order."""
        for child in doc.element.body:
            if child.tag == qn("w:p"):
                yield Paragraph(child, doc)
            elif child.tag == qn("w:tbl"):
                yield Table(child, doc)

    def _table_to_text(table) -> str:
        """Render a table as pipe-separated rows, deduplicating merged cells."""
        rows = []
        for row in table.rows:
            seen = set()
            cells = []
            for cell in row.cells:
                t = cell.text.strip()
                if t and t not in seen:
                    cells.append(t)
                    seen.add(t)
            if cells:
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    chunks = []
    buffer = ""
    for block in _iter_blocks(doc):
        if isinstance(block, Paragraph):
            text = block.text.strip()
        else:
            text = _table_to_text(block)
        if not text:
            continue
        if len(buffer) + len(text) > MAX_CHUNK_CHARS and buffer:
            chunks.append({"text": buffer.strip(), "source": str(path), "repo": repo, "chunk_type": "prose", "start_line": None})
            buffer = buffer[-OVERLAP_CHARS:] + "\n" + text
        else:
            buffer = (buffer + "\n" + text) if buffer else text
    if buffer.strip():
        chunks.append({"text": buffer.strip(), "source": str(path), "repo": repo, "chunk_type": "prose", "start_line": None})
    return chunks


EMPTY_PDF_LOG = Path(os.environ.get("EMPTY_PDF_LOG", "/tmp/empty_pdfs.log"))


def _log_empty_pdf(path: Path) -> None:
    try:
        with open(EMPTY_PDF_LOG, "a", encoding="utf-8") as f:
            f.write(str(path) + "\n")
    except Exception:
        pass


def _extract_text_pypdf(path: Path) -> str:
    """Fast text extraction using pypdf. Returns empty string on failure."""
    try:
        from pypdf import PdfReader
    except ImportError:
        return ""
    try:
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return ""


_AFFILIATION_WORDS = ("university", "department", "institute", "laboratory", "école", "universitat")

def _is_likely_paper(path: Path, text: str) -> bool:
    """Heuristic: Abstract + author signals (email/affiliation) or two-column layout.

    Deliberately does NOT use the folder name — a folder called 'Papers & Books'
    contains textbooks too, and routing those through Docling is very slow.
    Content signals are more reliable than path signals.
    """
    sample = text[:3000].lower()
    has_abstract = "abstract" in sample
    if not has_abstract:
        return False
    # Author signals: email address or known affiliation keywords
    has_authors = "@" in sample or any(w in sample for w in _AFFILIATION_WORDS)
    if has_authors:
        return True
    # Fallback: short average line length suggests two-column layout
    lines = [l for l in text[:3000].splitlines() if l.strip()]
    if lines and sum(len(l) for l in lines) / len(lines) < 65:
        return True
    return False


def _get_pdf_converter():
    """Return a cached DocumentConverter (loads models once per process).

    Set DOCLING_DEVICE=cuda to use GPU (faster, requires vLLM to be stopped).
    Defaults to CPU so it can run alongside vLLM.

    Set DOCLING_OCR=1 to enable OCR (for scanned PDFs). Slower but needed for
    image-only documents. Use with --file-list to target specific files.
    """
    from docling.document_converter import DocumentConverter, PdfFormatOption
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions, AcceleratorOptions, AcceleratorDevice
    device_str = os.environ.get("DOCLING_DEVICE", "cpu").lower()
    if device_str == "cuda":
        device = AcceleratorDevice.CUDA
    elif device_str == "auto":
        device = AcceleratorDevice.AUTO
    else:
        device = AcceleratorDevice.CPU
    do_ocr = os.environ.get("DOCLING_OCR", "").strip() == "1"
    pipeline_options = PdfPipelineOptions()
    pipeline_options.accelerator_options = AcceleratorOptions(device=device)
    pipeline_options.do_ocr = do_ocr
    return DocumentConverter(
        format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)}
    )


_pdf_converter = None  # None = not yet initialised; False = init failed permanently


def _chunk_pdf_text(text: str, path: Path, repo: str) -> list[dict]:
    """Chunk a PDF text string (already extracted) and tag source as path."""
    import tempfile
    # PDF extractors emit LONE UTF-16 surrogates from equation glyphs (e.g.
    # \ud835 = the high half of mathematical-italic letters, U+1D400-1D7FF in
    # physics papers). A lone surrogate is unencodable: the tempfile write
    # below raised UnicodeEncodeError and killed the whole batch, every night
    # (nightly sweep batch_00008, found 2026-07-24 via server_sweep.log).
    # Drop unpaired surrogates; well-formed text is unaffected.
    text = text.encode("utf-8", errors="ignore").decode("utf-8", errors="ignore")
    with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False, encoding="utf-8") as f:
        f.write(text)
        tmp = Path(f.name)
    try:
        chunks = _chunk_text(tmp, repo, chunk_type="prose")
    finally:
        tmp.unlink(missing_ok=True)
    for c in chunks:
        c["source"] = str(path)
    return chunks


def _chunk_pdf_docling(path: Path, repo: str) -> list[dict]:
    """Chunk PDF using Docling (layout-aware, handles two-column papers)."""
    global _pdf_converter
    if _pdf_converter is False:
        return []
    try:
        if _pdf_converter is None:
            _pdf_converter = _get_pdf_converter()
        result = _pdf_converter.convert(str(path))
        text = result.document.export_to_markdown()
        if not text.strip():
            _log_empty_pdf(path)
            return []
        chunks = _chunk_pdf_text(text, path, repo)
        if not chunks:
            _log_empty_pdf(path)
        return chunks
    except ImportError:
        logger.warning("Docling not installed — falling back to pypdf for all PDFs")
        _pdf_converter = False  # avoid retrying on every PDF
        return []
    except Exception as exc:
        logger.warning("PDF chunking failed for %s: %s", path, exc)
        return []


def _chunk_pdf(path: Path, repo: str) -> list[dict]:
    """Hybrid PDF chunker:
    - pypdf (fast) for manuals, books, reports
    - Docling (layout-aware) for papers and short/ambiguous PDFs
    - skip entirely if likely scanned (pypdf gets almost nothing, Docling can't OCR)

    Set DOCLING_DISABLE=1 to force pypdf-only mode (useful for bulk ingestion
    where speed matters and Docling can be run later on confirmed papers).
    """
    # Step 1: fast extraction with pypdf
    text = _extract_text_pypdf(path)

    MIN_PYPDF_CHARS = 2000   # below this, pypdf likely got only metadata/title
    MIN_DOCLING_CHARS = 200  # below this, PDF is likely scanned — Docling won't help without OCR

    # Step 2: if DOCLING_DISABLE is set, skip all Docling and use pypdf only.
    # Useful for bulk ingestion — run Docling later on confirmed papers.
    if os.environ.get("DOCLING_DISABLE", "").strip() == "1":
        if len(text.strip()) < MIN_DOCLING_CHARS:
            _log_empty_pdf(path)
            return []
        chunks = _chunk_pdf_text(text, path, repo)
        if chunks:
            return chunks
        _log_empty_pdf(path)
        return []

    # Step 3: if pypdf got almost nothing, it's almost certainly a scanned PDF.
    # Without OCR, Docling reads the same text layer and won't do better — skip it.
    # But when DOCLING_OCR=1, send these to Docling for OCR extraction.
    if len(text.strip()) < MIN_DOCLING_CHARS:
        if os.environ.get("DOCLING_OCR", "").strip() == "1":
            logger.debug("PDF likely scanned (pypdf got %d chars), sending to Docling OCR: %s", len(text.strip()), path.name)
            return _chunk_pdf_docling(path, repo)
        logger.debug("PDF likely scanned (pypdf got %d chars), skipping Docling: %s", len(text.strip()), path.name)
        _log_empty_pdf(path)
        return []

    # Step 4: if pypdf got a usable TEXT LAYER, extract it directly (FAST — skip
    # the slow Docling layout pass). PDF_TEXTLAYER_FAST=1 (bulk ingest) lowers the
    # bar to MIN_DOCLING_CHARS and includes papers: a born-digital PDF gains little
    # from Docling-WITHOUT-OCR (Docling reads the same text layer pypdf already
    # got, but ~50x slower on image-heavy pages). Truly-scanned PDFs (< the floor,
    # no text layer) fall through to Step 3's OCR path.
    fast_textlayer = os.environ.get("PDF_TEXTLAYER_FAST", "").strip() == "1"
    pypdf_floor = MIN_DOCLING_CHARS if fast_textlayer else MIN_PYPDF_CHARS
    if len(text.strip()) >= pypdf_floor and (fast_textlayer or not _is_likely_paper(path, text)):
        chunks = _chunk_pdf_text(text, path, repo)
        if chunks:
            logger.debug("PDF via pypdf (text layer): %s → %d chunks", path.name, len(chunks))
            return chunks

    # Step 5: paper detected or 200–2000 chars (short/ambiguous) → Docling for layout-aware extraction
    if len(text.strip()) < MIN_PYPDF_CHARS:
        reason = "insufficient text from pypdf"
    else:
        reason = "paper detected"
    logger.debug("PDF via Docling (%s): %s", reason, path.name)
    return _chunk_pdf_docling(path, repo)


# ── Dispatcher ─────────────────────────────────────────────────────────────────

def chunk_file(path: Path, repo: str) -> list[dict]:
    """Return chunks for a single file. Returns [] for unsupported types."""
    suffix = path.suffix.lower()
    if suffix == ".py":
        return _chunk_python(path, repo)
    if suffix == ".ipynb":
        return _chunk_notebook(path, repo)
    if suffix in (".md", ".txt", ".rst"):
        return _chunk_text(path, repo, chunk_type="prose")
    if suffix == ".pdf":
        return _chunk_pdf(path, repo)
    if suffix == ".pptx":
        return _chunk_pptx(path, repo)
    if suffix == ".docx":
        return _chunk_docx(path, repo)
    return []


# ── Helpers ────────────────────────────────────────────────────────────────────

def _sliding_window(text: str) -> list[str]:
    if len(text) <= MAX_CHUNK_CHARS:
        return [text]
    parts = []
    start = 0
    while start < len(text):
        end = start + MAX_CHUNK_CHARS
        parts.append(text[start:end])
        start += MAX_CHUNK_CHARS - OVERLAP_CHARS
    return parts
